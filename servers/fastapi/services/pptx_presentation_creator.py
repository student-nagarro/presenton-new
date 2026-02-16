import os
import re
import zipfile
from typing import List, Optional
from lxml import etree
from services.html_to_text_runs_service import (
    parse_html_text_to_text_runs as parse_inline_html_to_runs,
)

from pptx import Presentation
from pptx.shapes.autoshape import Shape
from pptx.slide import Slide
from pptx.text.text import _Paragraph, TextFrame, Font, _Run
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from lxml.etree import fromstring, tostring
from PIL import Image
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

from pptx.util import Pt
from pptx.dml.color import RGBColor

from models.pptx_models import (
    PptxAutoShapeBoxModel,
    PptxBoxShapeEnum,
    PptxConnectorModel,
    PptxFillModel,
    PptxFontModel,
    PptxParagraphModel,
    PptxPictureBoxModel,
    PptxPositionModel,
    PptxPresentationModel,
    PptxShadowModel,
    PptxSlideModel,
    PptxSpacingModel,
    PptxStrokeModel,
    PptxTextBoxModel,
    PptxTextRunModel,
)
from utils.download_helpers import download_files
from utils.image_utils import (
    clip_image,
    create_circle_image,
    fit_image,
    invert_image,
    round_image_corners,
    set_image_opacity,
)
import uuid

BLANK_SLIDE_LAYOUT = 6


class PptxPresentationCreator:
    def __init__(
        self,
        ppt_model: PptxPresentationModel,
        temp_dir: str,
        template_path: Optional[str] = None,
    ):
        self._temp_dir = temp_dir

        self._ppt_model = ppt_model
        self._slide_models = ppt_model.slides

        template_path = (template_path or "").strip()
        if template_path:
            template_path = os.path.abspath(os.path.expanduser(template_path))
            if not os.path.isfile(template_path):
                raise FileNotFoundError(
                    f"PPTX template not found: {template_path}"
                )
            self._ppt = Presentation(template_path)
            self._remove_all_slides()
        else:
            self._ppt = Presentation()
            self._ppt.slide_width = Pt(960)
            self._ppt.slide_height = Pt(540)

    def get_sub_element(self, parent, tagname, **kwargs):
        """Helper method to create XML elements"""
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

    def _remove_all_slides(self) -> None:
        # python-pptx has no public slide delete API; use the underlying XML list.
        sldIdLst = self._ppt.slides._sldIdLst  # type: ignore[attr-defined]
        for sldId in list(sldIdLst):
            rId = sldId.rId
            self._ppt.part.drop_rel(rId)
            sldIdLst.remove(sldId)

    async def fetch_network_assets(self):
        image_urls = []
        models_with_network_asset: List[PptxPictureBoxModel] = []

        if self._ppt_model.shapes:
            for each_shape in self._ppt_model.shapes:
                if isinstance(each_shape, PptxPictureBoxModel):
                    image_path = each_shape.picture.path
                    if image_path.startswith("http"):
                        if "app_data/" in image_path:
                            relative_path = image_path.split("app_data/")[1]
                            each_shape.picture.path = os.path.join(
                                "/app_data", relative_path
                            )
                            each_shape.picture.is_network = False
                            continue
                        image_urls.append(image_path)
                        models_with_network_asset.append(each_shape)

        for each_slide in self._slide_models:
            for each_shape in each_slide.shapes:
                if isinstance(each_shape, PptxPictureBoxModel):
                    image_path = each_shape.picture.path
                    if image_path.startswith("http"):
                        if "app_data" in image_path:
                            relative_path = image_path.split("app_data/")[1]
                            each_shape.picture.path = os.path.join(
                                "/app_data", relative_path
                            )
                            each_shape.picture.is_network = False
                            continue
                        image_urls.append(image_path)
                        models_with_network_asset.append(each_shape)

        if image_urls:
            image_paths = await download_files(image_urls, self._temp_dir)

            for each_shape, each_image_path in zip(
                models_with_network_asset, image_paths
            ):
                if each_image_path:
                    each_shape.picture.path = each_image_path
                    each_shape.picture.is_network = False

    async def create_ppt(self):
        await self.fetch_network_assets()

        self.set_theme_fonts()

        for slide_model in self._slide_models:
            # Adding global shapes to slide
            if self._ppt_model.shapes:
                slide_model.shapes.append(self._ppt_model.shapes)

            self.add_and_populate_slide(slide_model)

    def set_presentation_theme(self):
        slide_master = self._ppt.slide_master
        slide_master_part = slide_master.part

        theme_part = slide_master_part.part_related_by(RT.THEME)
        theme = fromstring(theme_part.blob)

        theme_colors = self._theme.colors.theme_color_mapping
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

        for color_name, hex_value in theme_colors.items():
            if color_name:
                color_element = theme.xpath(
                    f"a:themeElements/a:clrScheme/a:{color_name}/a:srgbClr",
                    namespaces=nsmap,
                )[0]
                color_element.set("val", hex_value.encode("utf-8"))

        theme_part._blob = tostring(theme)

    def set_theme_fonts(
        self,
        major_latin: str = "Equip Extended Medium",
        minor_latin: str = "Equip",
    ):
        slide_master = self._ppt.slide_master
        theme_part = slide_master.part.part_related_by(RT.THEME)
        theme = fromstring(theme_part.blob)

        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

        major_fonts = theme.xpath(
            "a:themeElements/a:fontScheme/a:majorFont/a:latin",
            namespaces=nsmap,
        )
        if major_fonts:
            major_fonts[0].set("typeface", major_latin)

        minor_fonts = theme.xpath(
            "a:themeElements/a:fontScheme/a:minorFont/a:latin",
            namespaces=nsmap,
        )
        if minor_fonts:
            minor_fonts[0].set("typeface", minor_latin)

        theme_part._blob = tostring(theme)

    def add_and_populate_slide(self, slide_model: PptxSlideModel):
        slide = self._ppt.slides.add_slide(self._ppt.slide_layouts[BLANK_SLIDE_LAYOUT])

        if slide_model.background:
            self.apply_fill_to_shape(slide.background, slide_model.background)

        if slide_model.note:
            slide.notes_slide.notes_text_frame.text = slide_model.note

        shapes = slide_model.shapes
        if shapes:
            slide_width_pt = float(self._ppt.slide_width.pt)
            slide_height_pt = float(self._ppt.slide_height.pt)
            # Ensure full-bleed background images are added first (bottom layer).
            full_bleed_backgrounds = []
            other_shapes = []
            for shape_model in shapes:
                if isinstance(shape_model, PptxPictureBoxModel):
                    pos = shape_model.position
                    if (
                        pos
                        and abs(pos.left) <= 1
                        and abs(pos.top) <= 1
                        and abs(pos.width - slide_width_pt) <= 1
                        and abs(pos.height - slide_height_pt) <= 1
                    ):
                        full_bleed_backgrounds.append(shape_model)
                        continue
                other_shapes.append(shape_model)
            shapes = full_bleed_backgrounds + other_shapes

        for shape_model in shapes:
            model_type = type(shape_model)

            if model_type is PptxPictureBoxModel:
                self.add_picture(slide, shape_model)

            elif model_type is PptxAutoShapeBoxModel:
                self.add_autoshape(slide, shape_model)

            elif model_type is PptxTextBoxModel:
                self.add_textbox(slide, shape_model)

            elif model_type is PptxConnectorModel:
                self.add_connector(slide, shape_model)

    def add_connector(self, slide: Slide, connector_model: PptxConnectorModel):
        if connector_model.thickness == 0:
            return
        connector_shape = slide.shapes.add_connector(
            connector_model.type, *connector_model.position.to_pt_xyxy()
        )
        connector_shape.line.width = Pt(connector_model.thickness)
        connector_shape.line.color.rgb = RGBColor.from_string(connector_model.color)
        self.set_fill_opacity(connector_shape, connector_model.opacity)

    def add_picture(self, slide: Slide, picture_model: PptxPictureBoxModel):
        image_path = picture_model.picture.path
        if (
            picture_model.clip
            or picture_model.border_radius
            or picture_model.invert
            or picture_model.opacity
            or picture_model.object_fit
            or picture_model.shape
        ):
            try:
                image = Image.open(image_path)
            except Exception:
                print(f"Could not open image: {image_path}")
                return

            image = image.convert("RGBA")
            # ? Applying border radius twice to support both clip and object fit
            if picture_model.border_radius:
                image = round_image_corners(image, picture_model.border_radius)
            if picture_model.object_fit:
                image = fit_image(
                    image,
                    picture_model.position.width,
                    picture_model.position.height,
                    picture_model.object_fit,
                )
            elif picture_model.clip:
                image = clip_image(
                    image,
                    picture_model.position.width,
                    picture_model.position.height,
                )
            if picture_model.border_radius:
                image = round_image_corners(image, picture_model.border_radius)
            if picture_model.shape == PptxBoxShapeEnum.CIRCLE:
                image = create_circle_image(image)
            if picture_model.invert:
                image = invert_image(image)
            if picture_model.opacity:
                image = set_image_opacity(image, picture_model.opacity)
            image_path = os.path.join(self._temp_dir, f"{uuid.uuid4()}.png")
            image.save(image_path)

        margined_position = self.get_margined_position(
            picture_model.position, picture_model.margin
        )

        slide.shapes.add_picture(image_path, *margined_position.to_pt_list())

    def add_autoshape(self, slide: Slide, autoshape_box_model: PptxAutoShapeBoxModel):
        position = autoshape_box_model.position
        if autoshape_box_model.margin:
            position = self.get_margined_position(position, autoshape_box_model.margin)

        autoshape = slide.shapes.add_shape(
            autoshape_box_model.type, *position.to_pt_list()
        )

        textbox = autoshape.text_frame
        textbox.auto_size = MSO_AUTO_SIZE.NONE
        textbox.word_wrap = autoshape_box_model.text_wrap

        self.apply_fill_to_shape(autoshape, autoshape_box_model.fill)
        self.apply_margin_to_text_box(textbox, autoshape_box_model.margin)
        self.apply_stroke_to_shape(autoshape, autoshape_box_model.stroke)
        self.apply_shadow_to_shape(autoshape, autoshape_box_model.shadow)
        self.apply_border_radius_to_shape(autoshape, autoshape_box_model.border_radius)

        if autoshape_box_model.paragraphs:
            self.add_paragraphs(textbox, autoshape_box_model.paragraphs)

    def add_textbox(self, slide: Slide, textbox_model: PptxTextBoxModel):
        position = textbox_model.position
        textbox_shape = slide.shapes.add_textbox(*position.to_pt_list())
        textbox_shape.width += Pt(2)

        textbox = textbox_shape.text_frame
        textbox.auto_size = MSO_AUTO_SIZE.NONE
        textbox.word_wrap = textbox_model.text_wrap

        self.apply_fill_to_shape(textbox_shape, textbox_model.fill)
        self.apply_margin_to_text_box(textbox, textbox_model.margin)
        self.add_paragraphs(textbox, textbox_model.paragraphs)

    def add_paragraphs(
        self, textbox: TextFrame, paragraph_models: List[PptxParagraphModel]
    ):
        for index, paragraph_model in enumerate(paragraph_models):
            paragraph = textbox.add_paragraph() if index > 0 else textbox.paragraphs[0]
            self.populate_paragraph(paragraph, paragraph_model)

    def populate_paragraph(
        self, paragraph: _Paragraph, paragraph_model: PptxParagraphModel
    ):
        if paragraph_model.list_level is not None:
            paragraph.level = paragraph_model.list_level

        list_type = paragraph_model.list_type or (
            "ul" if paragraph_model.list_level is not None else None
        )
        if list_type:
            self.apply_bullet_to_paragraph(
                paragraph, list_type, paragraph_model.list_item_index
            )

        list_indent = paragraph_model.list_indent
        list_hanging = paragraph_model.list_hanging
        if list_type:
            font_size = paragraph_model.font.size if paragraph_model.font else 16
            min_hanging = max(12, round(font_size * 0.9))
            desired_hanging = max(min_hanging, round(font_size * 1.2))
            if list_hanging is None or list_hanging < min_hanging:
                list_hanging = desired_hanging

            explicit_zero = list_indent == 0
            if list_indent is None:
                list_indent = round(font_size * 1.4)
            if not explicit_zero:
                min_indent = list_hanging + max(8, round(font_size * 0.4))
                if list_indent < min_indent:
                    list_indent = min_indent

        if list_indent is not None or list_hanging is not None:
            self.apply_list_indents(paragraph, list_indent, list_hanging)

        if paragraph_model.spacing:
            self.apply_spacing_to_paragraph(paragraph, paragraph_model.spacing)

        if paragraph_model.line_height:
            paragraph.line_spacing = paragraph_model.line_height

        if paragraph_model.alignment:
            paragraph.alignment = paragraph_model.alignment

        if paragraph_model.font:
            self.apply_font_to_paragraph(paragraph, paragraph_model.font)

        text_runs = []
        if paragraph_model.text:
            text_runs = self.parse_html_text_to_text_runs(
                paragraph_model.font, paragraph_model.text
            )
        elif paragraph_model.text_runs:
            text_runs = paragraph_model.text_runs

        if list_type and text_runs:
            self.strip_bullet_prefix_from_first_run(text_runs, list_type)

        for text_run_model in text_runs:
            text_run = paragraph.add_run()
            self.populate_text_run(text_run, text_run_model)

    def strip_bullet_prefix_from_first_run(
        self, text_runs: List[PptxTextRunModel], list_type: str
    ):
        if not text_runs:
            return

        first = text_runs[0]
        if not first.text:
            return

        original = first.text
        stripped = original

        if list_type == "ol":
            stripped = re.sub(r"^\s*\d+[\.\)]\s+", "", stripped, count=1)
        else:
            stripped = re.sub(r"^\s*[\u2022\u00B7]\s+", "", stripped, count=1)
            stripped = re.sub(r"^\s*-\s+", "", stripped, count=1)

        if stripped != original:
            first.text = stripped
            if first.text == "" and len(text_runs) > 1:
                text_runs.pop(0)

    def apply_bullet_to_paragraph(
        self, paragraph: _Paragraph, list_type: str, list_item_index: Optional[int]
    ):
        p_pr = paragraph._p.get_or_add_pPr()

        bu_none = p_pr.find(qn("a:buNone"))
        if bu_none is not None:
            p_pr.remove(bu_none)

        bu_char = p_pr.find(qn("a:buChar"))
        if bu_char is not None:
            p_pr.remove(bu_char)
        bu_auto = p_pr.find(qn("a:buAutoNum"))
        if bu_auto is not None:
            p_pr.remove(bu_auto)

        if list_type == "ol":
            bu = OxmlElement("a:buAutoNum")
            bu.set("type", "arabicPeriod")
            if list_item_index is not None:
                bu.set("startAt", str(list_item_index + 1))
            p_pr.append(bu)
        else:
            bu = OxmlElement("a:buChar")
            bu.set("char", "•")
            p_pr.append(bu)

    def apply_list_indents(
        self, paragraph: _Paragraph, indent: Optional[int], hanging: Optional[int]
    ):
        p_pr = paragraph._p.get_or_add_pPr()
        if indent is not None:
            emu = int(Pt(indent))
            p_pr.marL = emu
            p_pr.set("marL", str(emu))
        if hanging is not None:
            emu = int(Pt(-hanging))
            p_pr.indent = emu
            p_pr.set("indent", str(emu))

    def parse_html_text_to_text_runs(self, font: Optional[PptxFontModel], text: str):
        return parse_inline_html_to_runs(text, font)

    def populate_text_run(self, text_run: _Run, text_run_model: PptxTextRunModel):
        text_run.text = text_run_model.text
        if text_run_model.font:
            self.apply_font(text_run.font, text_run_model.font)

    def apply_border_radius_to_shape(self, shape: Shape, border_radius: Optional[int]):
        if not border_radius:
            return
        try:
            normalized_border_radius = Pt(border_radius) / min(
                shape.width, shape.height
            )
            shape.adjustments[0] = normalized_border_radius
        except Exception:
            print("Could not apply border radius.")

    def apply_fill_to_shape(self, shape: Shape, fill: Optional[PptxFillModel] = None):
        if not fill:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(fill.color)
            self.set_fill_opacity(shape.fill, fill.opacity)

    def apply_stroke_to_shape(
        self, shape: Shape, stroke: Optional[PptxStrokeModel] = None
    ):
        if not stroke or stroke.thickness == 0:
            shape.line.fill.background()
        else:
            shape.line.fill.solid()
            shape.line.fill.fore_color.rgb = RGBColor.from_string(stroke.color)
            shape.line.width = Pt(stroke.thickness)
            self.set_fill_opacity(shape.line.fill, stroke.opacity)

    def apply_shadow_to_shape(
        self, shape: Shape, shadow: Optional[PptxShadowModel] = None
    ):
        # Access the XML for the shape
        sp_element = shape._element
        sp_pr = sp_element.xpath("p:spPr")[0]  # Shape properties XML element

        nsmap = sp_pr.nsmap

        # # Remove existing shadow effects if present
        effect_list = sp_pr.find("a:effectLst", namespaces=nsmap)
        if effect_list:
            old_outer_shadow = effect_list.find("a:outerShdw")
            if old_outer_shadow:
                effect_list.remove(
                    old_outer_shadow, namespaces=nsmap
                )  # Remove the old shadow
            old_inner_shadow = effect_list.find("a:innerShdw")
            if old_inner_shadow:
                effect_list.remove(
                    old_inner_shadow, namespaces=nsmap
                )  # Remove the old shadow
            old_prst_shadow = effect_list.find("a:prstShdw")
            if old_prst_shadow:
                effect_list.remove(
                    old_prst_shadow, namespaces=nsmap
                )  # Remove the old shadow

        if not effect_list:
            effect_list = etree.SubElement(
                sp_pr, f"{{{nsmap['a']}}}effectLst", nsmap=nsmap
            )

        if shadow is None:
            # Apply shadow with zero values when shadow is None
            outer_shadow = etree.SubElement(
                effect_list,
                f"{{{nsmap['a']}}}outerShdw",
                {
                    "blurRad": "0",
                    "dist": "0",
                    "dir": "0",
                },
                nsmap=nsmap,
            )
            color_element = etree.SubElement(
                outer_shadow,
                f"{{{nsmap['a']}}}srgbClr",
                {"val": "000000"},
                nsmap=nsmap,
            )
            etree.SubElement(
                color_element,
                f"{{{nsmap['a']}}}alpha",
                {"val": "0"},
                nsmap=nsmap,
            )
        else:
            # Apply the provided shadow
            # dir expects 60000ths of a degree in OOXML
            angle_dir = (
                int(round((shadow.angle % 360) * 60000))
                if shadow.angle is not None
                else 0
            )
            outer_shadow = etree.SubElement(
                effect_list,
                f"{{{nsmap['a']}}}outerShdw",
                {
                    "blurRad": f"{Pt(shadow.radius)}",
                    "dir": f"{angle_dir}",
                    "dist": f"{Pt(shadow.offset)}",
                    "rotWithShape": "0",
                },
                nsmap=nsmap,
            )
            color_element = etree.SubElement(
                outer_shadow,
                f"{{{nsmap['a']}}}srgbClr",
                {"val": f"{shadow.color}"},
                nsmap=nsmap,
            )
            etree.SubElement(
                color_element,
                f"{{{nsmap['a']}}}alpha",
                {"val": f"{int(shadow.opacity * 100000)}"},
                nsmap=nsmap,
            )

    def set_fill_opacity(self, fill, opacity):
        if opacity is None or opacity >= 1.0:
            return

        alpha = int((opacity) * 100000)

        try:
            ts = fill._xPr.solidFill
            sF = ts.get_or_change_to_srgbClr()
            self.get_sub_element(sF, "a:alpha", val=str(alpha))
        except Exception as e:
            print(f"Could not set fill opacity: {e}")

    def get_margined_position(
        self, position: PptxPositionModel, margin: Optional[PptxSpacingModel]
    ) -> PptxPositionModel:
        if not margin:
            return position

        left = position.left + margin.left
        top = position.top + margin.top
        width = max(position.width - margin.left - margin.right, 0)
        height = max(position.height - margin.top - margin.bottom, 0)

        return PptxPositionModel(left=left, top=top, width=width, height=height)

    def apply_margin_to_text_box(
        self, text_frame: TextFrame, margin: Optional[PptxSpacingModel]
    ) -> PptxPositionModel:
        text_frame.margin_left = Pt(margin.left if margin else 0)
        text_frame.margin_right = Pt(margin.right if margin else 0)
        text_frame.margin_top = Pt(margin.top if margin else 0)
        text_frame.margin_bottom = Pt(margin.bottom if margin else 0)

    def apply_spacing_to_paragraph(
        self, paragraph: _Paragraph, spacing: PptxSpacingModel
    ):
        paragraph.space_before = Pt(spacing.top)
        paragraph.space_after = Pt(spacing.bottom)

    def apply_font_to_paragraph(self, paragraph: _Paragraph, font: PptxFontModel):
        self.apply_font(paragraph.font, font)

    def apply_font(self, font: Font, font_model: PptxFontModel):
        font.name = font_model.name
        font.color.rgb = RGBColor.from_string(font_model.color)
        font.italic = font_model.italic
        font.size = Pt(font_model.size)
        font.bold = font_model.font_weight >= 600
        if font.name == "Equip Extended ExtraBold":
            font.bold = True
        if font.name == "Equip Extended Light" and font.bold:
            font.bold = False
        if font.name == "Equip Medium" and (font.bold or font.italic):
            font.name = "Equip"
        if font_model.underline is not None:
            font.underline = bool(font_model.underline)
        if font_model.strike is not None:
            self.apply_strike_to_font(font, font_model.strike)

    def apply_strike_to_font(self, font: Font, strike: Optional[bool]):
        try:
            rPr = font._element
            if strike is True:
                rPr.set("strike", "sngStrike")
            elif strike is False:
                rPr.set("strike", "noStrike")
        except Exception as e:
            print(f"Could not apply strikethrough: {e}")

    def save(self, path: str):
        self._ppt.save(path)
        self._force_last_view_normal(path)

    def _force_last_view_normal(self, path: str):
        temp_path = f"{path}.tmp"
        try:
            with zipfile.ZipFile(path, "r") as zin:
                names = set(zin.namelist())
                view_props = None
                if "ppt/viewProps.xml" in names:
                    view_props = zin.read("ppt/viewProps.xml")

                if view_props:
                    try:
                        root = fromstring(view_props)
                        root.set("lastView", "sldView")
                        view_props = tostring(
                            root, xml_declaration=True, encoding="UTF-8", standalone="yes"
                        )
                    except Exception:
                        view_props = re.sub(
                            r'lastView="[^"]+"',
                            'lastView="sldView"',
                            view_props.decode("utf-8", errors="ignore"),
                        ).encode("utf-8")
                else:
                    view_props = (
                        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                        '<p:viewPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
                        'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
                        'lastView="sldView"/>'
                    ).encode("utf-8")

                with zipfile.ZipFile(temp_path, "w") as zout:
                    for item in zin.infolist():
                        if item.filename == "ppt/viewProps.xml":
                            continue
                        data = zin.read(item.filename)
                        zout.writestr(item, data)
                    zout.writestr("ppt/viewProps.xml", view_props)
            os.replace(temp_path, path)
        except Exception as exc:
            try:
                if os.path.exists(temp_path):
                    os.remove(temp_path)
            except Exception:
                pass
            print(f"Could not set lastView to sldView: {exc}")
